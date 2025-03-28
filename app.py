import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO
import zipfile
import os
import subprocess
import tempfile
import shutil

# Configuración inicial de la app de Streamlit
title = "🎓 Generador de Certificados"
st.set_page_config(page_title=title, layout="centered")
st.title(title)

# Función que genera un certificado a partir de un nombre y una plantilla
def generate_certificate(name, template_bytes):
    prs = Presentation(BytesIO(template_bytes))
    output = BytesIO()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{NOMBRE}}" in run.text:
                            run.text = run.text.replace("{{NOMBRE}}", name)

    prs.save(output)
    output.seek(0)
    return output

# Función para convertir PPTX a PDF usando LibreOffice (si está disponible)
def convert_to_pdf(pptx_bytes, output_filename):
    libreoffice_cmd = "soffice"

    if not shutil.which(libreoffice_cmd):
        raise RuntimeError("❌ La conversión a PDF no está disponible en este entorno (se requiere LibreOffice).")

    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "temp_cert.pptx")
        pdf_dir = tmpdir

        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes.read())

        subprocess.run([
            libreoffice_cmd, "--headless", "--convert-to", "pdf", pptx_path, "--outdir", pdf_dir
        ], check=True)

        pdf_path = pptx_path.replace(".pptx", ".pdf")
        with open(pdf_path, "rb") as f:
            return f.read()

# Carga de archivos: plantilla PowerPoint y archivo Excel
uploaded_template = st.file_uploader(
    "Sube la plantilla PowerPoint (.pptx) con {{NOMBRE}}", type="pptx"
)
uploaded_excel = st.file_uploader(
    "Sube el archivo Excel con los participantes", type="xlsx"
)

# Determina opciones disponibles para el formato de salida
disponible_pdf = shutil.which("soffice") is not None
formatos_disponibles = ["PPTX"]
if disponible_pdf:
    formatos_disponibles.append("PDF")

# Selección de formato de salida
default_format = "PPTX"
output_format = st.selectbox("Selecciona el formato de salida", formatos_disponibles, index=formatos_disponibles.index(default_format))

# Verifica que ambos archivos hayan sido subidos
if uploaded_template and uploaded_excel:
    df = pd.read_excel(uploaded_excel)
    nombre_columna = st.selectbox("Selecciona la columna con los nombres", df.columns)

    if st.button("✅ Generar certificados"):
        zip_buffer = BytesIO()
        template_bytes = uploaded_template.read()

        try:
            with zipfile.ZipFile(zip_buffer, 'w') as zipf:
                for name in df[nombre_columna]:
                    cert = generate_certificate(name, template_bytes)
                    filename_base = f"Certificado_{name.replace(' ', '_')}"

                    if output_format == "PDF":
                        try:
                            pdf_content = convert_to_pdf(cert, filename_base + ".pdf")
                            zipf.writestr(filename_base + ".pdf", pdf_content)
                        except Exception as e:
                            st.error("❌ Error al convertir a PDF: " + str(e))
                            st.stop()
                    else:
                        zipf.writestr(filename_base + ".pptx", cert.read())

            zip_buffer.seek(0)
            st.success("✨ Certificados generados correctamente")
            st.download_button(
                label="📦 Descargar ZIP de certificados",
                data=zip_buffer,
                file_name="certificados.zip",
                mime="application/zip"
            )
        except Exception as e:
            st.error(f"❌ Ocurrió un error inesperado: {e}")
